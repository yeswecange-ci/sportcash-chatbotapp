from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── Palette TotalEnergie ────────────────────────────────────────────────────
RED       = RGBColor(0xE3, 0x00, 0x1E)   # rouge TotalEnergies
DARK_GRAY = RGBColor(0x1E, 0x1E, 0x2E)   # fond sombre
MID_GRAY  = RGBColor(0x2D, 0x2D, 0x3F)   # fond carte
LIGHT_BG  = RGBColor(0xF5, 0xF5, 0xF7)   # fond clair
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT    = RGBColor(0xFF, 0xC1, 0x07)   # jaune accent
GREEN     = RGBColor(0x28, 0xA7, 0x45)
BLUE      = RGBColor(0x0D, 0x6E, 0xFD)
TEXT_DARK = RGBColor(0x1E, 0x1E, 0x2E)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # blank layout

# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=Pt(14), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def screenshot_placeholder(slide, left, top, width, height, label="Capture d'écran"):
    """Cadre grisé avec croix et texte centré pour indiquer où coller la capture."""
    # fond
    add_rect(slide, left, top, width, height, fill_color=RGBColor(0xE8,0xE8,0xF0), line_color=RGBColor(0xAA,0xAA,0xBB), line_width=Pt(1.5))
    # croix diagonale (2 rectangles fins à 45°)
    cx = left + width  // 2
    cy = top  + height // 2
    diag_len = min(width, height) * 0.35
    # ligne 1 — fin rectangle horizontal (symbolique)
    add_rect(slide, left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Pt(1), fill_color=RGBColor(0xCC,0xCC,0xDD))
    add_rect(slide, left + Inches(0.1), top + height - Inches(0.15), width - Inches(0.2), Pt(1), fill_color=RGBColor(0xCC,0xCC,0xDD))
    # icône caméra textuelle + label
    add_text(slide, "📷", left, top + height//2 - Inches(0.45), width, Inches(0.5),
             font_size=Pt(28), align=PP_ALIGN.CENTER, color=RGBColor(0x88,0x88,0x99))
    add_text(slide, label, left, top + height//2 + Inches(0.1), width, Inches(0.4),
             font_size=Pt(11), align=PP_ALIGN.CENTER, color=RGBColor(0x66,0x66,0x88), italic=True)

def bullet_list(slide, items, left, top, width, height, icon="▸"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{icon}  {item}"
        run.font.size = Pt(13)
        run.font.color.rgb = TEXT_DARK

def badge(slide, text, left, top, width=Inches(1.8), height=Inches(0.35), color=RED):
    add_rect(slide, left, top, width, height, fill_color=color)
    add_text(slide, text, left, top, width, height,
             font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITRE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# fond sombre
add_rect(slide, 0, 0, W, H, fill_color=DARK_GRAY)

# bande rouge gauche
add_rect(slide, 0, 0, Inches(0.35), H, fill_color=RED)

# bande rouge basse
add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), fill_color=RED)

# Logo / nom app (simulé en texte)
add_text(slide, "WhatsCRM", Inches(0.6), Inches(1.2), Inches(6), Inches(1.2),
         font_size=Pt(64), bold=True, color=WHITE)
add_text(slide, "by YesWeCange × TotalEnergies", Inches(0.6), Inches(2.4), Inches(8), Inches(0.6),
         font_size=Pt(20), color=ACCENT)

# séparateur
add_rect(slide, Inches(0.6), Inches(3.1), Inches(5), Inches(0.04), fill_color=RED)

# sous-titre
add_text(slide,
         "Plateforme de support client & marketing WhatsApp\n"
         "avec chatbot automatisé, campagnes en masse et gamification",
         Inches(0.6), Inches(3.3), Inches(9), Inches(1.2),
         font_size=Pt(16), color=RGBColor(0xCC,0xCC,0xDD))

# Badges stack
badge(slide, "Laravel 12",   Inches(0.6),  Inches(5.2), color=RED)
badge(slide, "WhatsApp",     Inches(2.55), Inches(5.2), color=GREEN)
badge(slide, "Twilio",       Inches(4.5),  Inches(5.2), color=BLUE)
badge(slide, "Chatwoot",     Inches(6.45), Inches(5.2), color=RGBColor(0x14,0x7E,0xFF))
badge(slide, "Docker",       Inches(8.4),  Inches(5.2), color=MID_GRAY)

# Date
add_text(slide, "Présentation — Mars 2026", Inches(0.6), Inches(6.3), Inches(5), Inches(0.4),
         font_size=Pt(12), color=RGBColor(0x88,0x88,0x99), italic=True)

# image placeholder droite
screenshot_placeholder(slide, Inches(8.8), Inches(1.0), Inches(4.1), Inches(5.2),
                        "Logo / Écran d'accueil")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_color=LIGHT_BG)
add_rect(slide, 0, 0, W, Inches(1.1), fill_color=DARK_GRAY)
add_rect(slide, 0, Inches(1.08), W, Inches(0.04), fill_color=RED)

add_text(slide, "Sommaire", Inches(0.5), Inches(0.2), Inches(10), Inches(0.8),
         font_size=Pt(32), bold=True, color=WHITE)

modules = [
    ("01", "Accès & Connexion",           "Identifiants, URL et rôles utilisateurs"),
    ("02", "Dashboard",                    "Vue d'ensemble et indicateurs en temps réel"),
    ("03", "Conversations",               "Gestion des échanges WhatsApp agent ↔ client"),
    ("04", "Campagnes Marketing",         "Envoi en masse, templates et planification"),
    ("05", "Bot Tracking",                "Analyse du parcours chatbot et données clients"),
    ("06", "Gamification",               "Jeux, quiz et classements via WhatsApp"),
    ("07", "Agents & Équipes",           "Gestion des utilisateurs et attributions"),
    ("08", "Statistiques & Rapports",    "Métriques journalières et performances"),
]

cols = 2
rows = 4
col_w = Inches(6.2)
row_h = Inches(1.25)
start_left = Inches(0.4)
start_top  = Inches(1.3)

for i, (num, title, desc) in enumerate(modules):
    col = i % cols
    row = i // cols
    left = start_left + col * col_w
    top  = start_top  + row * row_h

    add_rect(slide, left, top, col_w - Inches(0.2), row_h - Inches(0.1),
             fill_color=WHITE, line_color=RGBColor(0xDD,0xDD,0xEE), line_width=Pt(1))
    add_rect(slide, left, top, Inches(0.45), row_h - Inches(0.1), fill_color=RED)
    add_text(slide, num, left, top, Inches(0.45), row_h - Inches(0.1),
             font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, left + Inches(0.55), top + Inches(0.12),
             col_w - Inches(0.85), Inches(0.45),
             font_size=Pt(14), bold=True, color=TEXT_DARK)
    add_text(slide, desc, left + Inches(0.55), top + Inches(0.58),
             col_w - Inches(0.85), Inches(0.5),
             font_size=Pt(11), color=RGBColor(0x55,0x55,0x77))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ACCÈS & CONNEXION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_color=LIGHT_BG)
add_rect(slide, 0, 0, W, Inches(1.1), fill_color=DARK_GRAY)
add_rect(slide, 0, Inches(1.08), W, Inches(0.04), fill_color=RED)

add_text(slide, "01 — Accès & Connexion", Inches(0.5), Inches(0.2), Inches(10), Inches(0.75),
         font_size=Pt(28), bold=True, color=WHITE)

# Colonne gauche — infos accès
add_rect(slide, Inches(0.4), Inches(1.3), Inches(5.5), Inches(5.8),
         fill_color=WHITE, line_color=RGBColor(0xDD,0xDD,0xEE), line_width=Pt(1))

add_text(slide, "Informations de connexion", Inches(0.6), Inches(1.45), Inches(5), Inches(0.5),
         font_size=Pt(16), bold=True, color=TEXT_DARK)
add_rect(slide, Inches(0.6), Inches(1.95), Inches(4.8), Inches(0.03), fill_color=RGBColor(0xEE,0xEE,0xEE))

rows_data = [
    ("🌐  URL Application",    "http://localhost:8080\n(ou domaine production)"),
    ("👤  Compte Admin",       "admin@totalenergies.ci\nMot de passe : [à définir]"),
    ("👤  Compte Agent",       "agent@totalenergies.ci\nMot de passe : [à définir]"),
    ("🔑  Chatwoot",           "URL : [CHATWOOT_BASE_URL]\nToken API : [CHATWOOT_API_TOKEN]"),
    ("📱  WhatsApp (Twilio)",  "Numéro : [TWILIO_FROM_NUMBER]\nSID : [TWILIO_SID]"),
]
top_cur = Inches(2.1)
for label, val in rows_data:
    add_text(slide, label, Inches(0.6), top_cur, Inches(2.2), Inches(0.55),
             font_size=Pt(11), bold=True, color=RED)
    add_text(slide, val, Inches(2.9), top_cur, Inches(2.8), Inches(0.6),
             font_size=Pt(11), color=TEXT_DARK)
    top_cur += Inches(0.78)

# Rôles
add_text(slide, "Rôles disponibles", Inches(0.6), Inches(6.15), Inches(5), Inches(0.4),
         font_size=Pt(13), bold=True, color=TEXT_DARK)
badge(slide, "super_admin", Inches(0.6),  Inches(6.6), color=RED)
badge(slide, "admin",       Inches(2.55), Inches(6.6), color=BLUE)
badge(slide, "agent",       Inches(4.15), Inches(6.6), color=GREEN)

# Colonne droite — screenshot
screenshot_placeholder(slide, Inches(6.3), Inches(1.3), Inches(6.6), Inches(5.8),
                        "Capture : Page de connexion")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — DASHBOARD
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_color=LIGHT_BG)
add_rect(slide, 0, 0, W, Inches(1.1), fill_color=DARK_GRAY)
add_rect(slide, 0, Inches(1.08), W, Inches(0.04), fill_color=RED)

add_text(slide, "02 — Dashboard", Inches(0.5), Inches(0.2), Inches(10), Inches(0.75),
         font_size=Pt(28), bold=True, color=WHITE)

# KPI cards (4)
kpis = [
    ("💬", "Conversations", "actives / totales"),
    ("📨", "Campagnes",    "envoyées ce mois"),
    ("🤖", "Bot Sessions", "journalières"),
    ("👥", "Agents",       "disponibles"),
]
card_w = Inches(2.9)
card_h = Inches(1.6)
for i, (ico, title, sub) in enumerate(kpis):
    lft = Inches(0.4) + i * (card_w + Inches(0.22))
    add_rect(slide, lft, Inches(1.25), card_w, card_h,
             fill_color=WHITE, line_color=RGBColor(0xDD,0xDD,0xEE), line_width=Pt(1))
    add_rect(slide, lft, Inches(1.25), Inches(0.07), card_h, fill_color=RED)
    add_text(slide, ico,   lft + Inches(0.2), Inches(1.35), card_w, Inches(0.55), font_size=Pt(24), color=TEXT_DARK)
    add_text(slide, title, lft + Inches(0.2), Inches(1.85), card_w - Inches(0.3), Inches(0.4),
             font_size=Pt(13), bold=True, color=TEXT_DARK)
    add_text(slide, sub,   lft + Inches(0.2), Inches(2.2),  card_w - Inches(0.3), Inches(0.35),
             font_size=Pt(10), color=RGBColor(0x77,0x77,0x88))

# grand screenshot dashboard
screenshot_placeholder(slide, Inches(0.4), Inches(3.05), Inches(12.5), Inches(4.0),
                        "Capture : Vue Dashboard complet (admin)")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CONVERSATIONS (description)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_color=LIGHT_BG)
add_rect(slide, 0, 0, W, Inches(1.1), fill_color=DARK_GRAY)
add_rect(slide, 0, Inches(1.08), W, Inches(0.04), fill_color=RED)

add_text(slide, "03 — Conversations WhatsApp", Inches(0.5), Inches(0.2), Inches(11), Inches(0.75),
         font_size=Pt(28), bold=True, color=WHITE)

add_rect(slide, Inches(0.4), Inches(1.25), Inches(5.5), Inches(5.9),
         fill_color=WHITE, line_color=RGBColor(0xDD,0xDD,0xEE), line_width=Pt(1))

add_text(slide, "Fonctionnement", Inches(0.6), Inches(1.4), Inches(5), Inches(0.5),
         font_size=Pt(15), bold=True, color=TEXT_DARK)

bullet_list(slide, [
    "Handoff automatique : Twilio Studio → Agent humain",
    "Fil de conversation temps réel (polling 4 secondes)",
    "Envoi de messages libres et templates WhatsApp",
    "Changement de statut : Actif / Transféré / Terminé / Abandonné",
    "Assignation à un agent ou une équipe",
    "Étiquettes (labels) pour catégoriser les échanges",
    "Réponses rapides (canned responses)",
    "Synchronisation bidirectionnelle avec Chatwoot",
], Inches(0.6), Inches(2.0), Inches(5.0), Inches(4.5))

screenshot_placeholder(slide, Inches(6.3), Inches(1.25), Inches(6.6), Inches(5.9),
                        "Capture : Interface conversations")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CAMPAGNES
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_color=LIGHT_BG)
add_rect(slide, 0, 0, W, Inches(1.1), fill_color=DARK_GRAY)
add_rect(slide, 0, Inches(1.08), W, Inches(0.04), fill_color=RED)

add_text(slide, "04 — Campagnes Marketing", Inches(0.5), Inches(0.2), Inches(11), Inches(0.75),
         font_size=Pt(28), bold=True, color=WHITE)

# Flux en étapes
steps = [
    ("1", "Créer",     "Nommer la campagne\net choisir le template"),
    ("2", "Contacts",  "Importer CSV\nou sync Chatwoot"),
    ("3", "Planifier", "Envoi immédiat\nou différé"),
    ("4", "Envoyer",   "Traitement asynchrone\nvia queue"),
    ("5", "Suivre",    "Livré / Échoué\nStatistiques"),
]
sw = Inches(2.3)
for i, (num, title, desc) in enumerate(steps):
    lft = Inches(0.3) + i * (sw + Inches(0.05))
    add_rect(slide, lft, Inches(1.25), sw, Inches(2.2),
             fill_color=WHITE, line_color=RGBColor(0xDD,0xDD,0xEE), line_width=Pt(1))
    add_rect(slide, lft, Inches(1.25), sw, Inches(0.4), fill_color=RED)
    add_text(slide, f"Étape {num}", lft, Inches(1.25), sw, Inches(0.4),
             font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, lft, Inches(1.7), sw, Inches(0.45),
             font_size=Pt(14), bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, desc, lft + Inches(0.1), Inches(2.2), sw - Inches(0.2), Inches(0.9),
             font_size=Pt(11), color=RGBColor(0x55,0x55,0x77), align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(slide, "→", lft + sw, Inches(1.9), Inches(0.12), Inches(0.4),
                 font_size=Pt(18), bold=True, color=RED, align=PP_ALIGN.CENTER)

screenshot_placeholder(slide, Inches(0.4), Inches(3.65), Inches(5.8), Inches(3.5),
                        "Capture : Création de campagne")
screenshot_placeholder(slide, Inches(6.6), Inches(3.65), Inches(6.3), Inches(3.5),
                        "Capture : Statistiques campagne")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — BOT TRACKING
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_color=LIGHT_BG)
add_rect(slide, 0, 0, W, Inches(1.1), fill_color=DARK_GRAY)
add_rect(slide, 0, Inches(1.08), W, Inches(0.04), fill_color=RED)

add_text(slide, "05 — Bot Tracking & Analytiques", Inches(0.5), Inches(0.2), Inches(11), Inches(0.75),
         font_size=Pt(28), bold=True, color=WHITE)

add_rect(slide, Inches(0.4), Inches(1.25), Inches(5.5), Inches(5.9),
         fill_color=WHITE, line_color=RGBColor(0xDD,0xDD,0xEE), line_width=Pt(1))

add_text(slide, "Ce que le module capture", Inches(0.6), Inches(1.4), Inches(5), Inches(0.5),
         font_size=Pt(15), bold=True, color=TEXT_DARK)

bullet_list(slide, [
    "Parcours complet de l'utilisateur dans le bot (menu_path)",
    "Choix de menus et saisies libres",
    "Données client : VIN, email, nom de profil, téléphone",
    "Statistiques journalières agrégées par menu et événement",
    "Détection des conversations abandonnées ou transférées",
    "Tableau de bord analytique avec filtres par date",
    "Export des données clients capturées",
], Inches(0.6), Inches(2.0), Inches(5.0), Inches(4.5))

screenshot_placeholder(slide, Inches(6.3), Inches(1.25), Inches(6.6), Inches(5.9),
                        "Capture : Bot Tracking / Tableau de bord analytique")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — GAMIFICATION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_color=LIGHT_BG)
add_rect(slide, 0, 0, W, Inches(1.1), fill_color=DARK_GRAY)
add_rect(slide, 0, Inches(1.08), W, Inches(0.04), fill_color=RED)

add_text(slide, "06 — Gamification & Quiz WhatsApp", Inches(0.5), Inches(0.2), Inches(11), Inches(0.75),
         font_size=Pt(28), bold=True, color=WHITE)

# 3 cartes fonctionnalités
feats = [
    ("🎮", "Créer un Jeu",
     "Définir les questions,\nréponses et durée du jeu.\nGénération automatique\ndu flow Twilio Studio."),
    ("📊", "Suivi en direct",
     "Participation, réponses,\nscores par joueur.\nTableau de bord temps réel."),
    ("🏆", "Classement",
     "Classement des meilleurs\njoueurs. Export des données.\nEngagement client gamifié."),
]
cw = Inches(3.9)
for i, (ico, title, desc) in enumerate(feats):
    lft = Inches(0.4) + i * (cw + Inches(0.22))
    add_rect(slide, lft, Inches(1.25), cw, Inches(3.1),
             fill_color=WHITE, line_color=RGBColor(0xDD,0xDD,0xEE), line_width=Pt(1))
    add_rect(slide, lft, Inches(1.25), cw, Inches(0.06), fill_color=RED)
    add_text(slide, ico,   lft, Inches(1.4),  cw, Inches(0.7), font_size=Pt(30), align=PP_ALIGN.CENTER, color=TEXT_DARK)
    add_text(slide, title, lft, Inches(2.1),  cw, Inches(0.5), font_size=Pt(15), bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, desc,  lft + Inches(0.15), Inches(2.65), cw - Inches(0.3), Inches(1.5),
             font_size=Pt(12), color=RGBColor(0x55,0x55,0x77), align=PP_ALIGN.CENTER)

screenshot_placeholder(slide, Inches(0.4), Inches(4.55), Inches(5.8), Inches(2.7),
                        "Capture : Interface création de jeu")
screenshot_placeholder(slide, Inches(6.6), Inches(4.55), Inches(6.3), Inches(2.7),
                        "Capture : Classement / Leaderboard")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — AGENTS & ÉQUIPES
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_color=LIGHT_BG)
add_rect(slide, 0, 0, W, Inches(1.1), fill_color=DARK_GRAY)
add_rect(slide, 0, Inches(1.08), W, Inches(0.04), fill_color=RED)

add_text(slide, "07 — Agents & Équipes", Inches(0.5), Inches(0.2), Inches(11), Inches(0.75),
         font_size=Pt(28), bold=True, color=WHITE)

add_rect(slide, Inches(0.4), Inches(1.25), Inches(5.5), Inches(5.9),
         fill_color=WHITE, line_color=RGBColor(0xDD,0xDD,0xEE), line_width=Pt(1))

add_text(slide, "Fonctionnalités RH & Support", Inches(0.6), Inches(1.4), Inches(5), Inches(0.5),
         font_size=Pt(15), bold=True, color=TEXT_DARK)

bullet_list(slide, [
    "Création et gestion des comptes agents",
    "Organisation en équipes avec membres assignés",
    "Statut de disponibilité agent (en ligne / hors ligne)",
    "Attribution automatique des conversations par équipe",
    "Règles d'auto-assignation configurables",
    "Gestion des tokens Chatwoot par agent",
    "Rôles hiérarchiques : super_admin > admin > agent",
    "Tableau de bord des performances par agent",
], Inches(0.6), Inches(2.0), Inches(5.0), Inches(4.5))

screenshot_placeholder(slide, Inches(6.3), Inches(1.25), Inches(6.6), Inches(5.9),
                        "Capture : Gestion agents / équipes")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — STATISTIQUES
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_color=LIGHT_BG)
add_rect(slide, 0, 0, W, Inches(1.1), fill_color=DARK_GRAY)
add_rect(slide, 0, Inches(1.08), W, Inches(0.04), fill_color=RED)

add_text(slide, "08 — Statistiques & Rapports", Inches(0.5), Inches(0.2), Inches(11), Inches(0.75),
         font_size=Pt(28), bold=True, color=WHITE)

metrics = [
    ("📈", "Conversations",  "Total, actives, terminées,\nabandon rate par jour"),
    ("📬", "Campagnes",      "Taux de livraison, échecs,\ntemps d'envoi moyen"),
    ("🤖", "Bot",            "Menus les plus utilisés,\nparcours populaires"),
    ("👤", "Agents",         "Conversations traitées,\ntemps de réponse moyen"),
    ("📅", "Journalier",     "Agrégats quotidiens\nautonomes (cron)"),
    ("📤", "Export",         "Données exportables\nau format CSV"),
]
mw = Inches(3.9)
mh = Inches(1.8)
for i, (ico, title, desc) in enumerate(metrics):
    col = i % 3
    row = i // 3
    lft = Inches(0.4) + col * (mw + Inches(0.22))
    top = Inches(1.3) + row * (mh + Inches(0.15))
    add_rect(slide, lft, top, mw, mh, fill_color=WHITE,
             line_color=RGBColor(0xDD,0xDD,0xEE), line_width=Pt(1))
    add_rect(slide, lft, top, Inches(0.06), mh, fill_color=RED)
    add_text(slide, ico + "  " + title, lft + Inches(0.15), top + Inches(0.15),
             mw - Inches(0.2), Inches(0.5), font_size=Pt(14), bold=True, color=TEXT_DARK)
    add_text(slide, desc, lft + Inches(0.15), top + Inches(0.7),
             mw - Inches(0.2), Inches(0.9), font_size=Pt(12), color=RGBColor(0x55,0x55,0x77))

screenshot_placeholder(slide, Inches(0.4), Inches(5.25), Inches(12.5), Inches(2.0),
                        "Capture : Page Statistiques complète")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — ARCHITECTURE TECHNIQUE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_color=LIGHT_BG)
add_rect(slide, 0, 0, W, Inches(1.1), fill_color=DARK_GRAY)
add_rect(slide, 0, Inches(1.08), W, Inches(0.04), fill_color=RED)

add_text(slide, "Architecture Technique", Inches(0.5), Inches(0.2), Inches(11), Inches(0.75),
         font_size=Pt(28), bold=True, color=WHITE)

# Blocs architecture
arch_items = [
    (Inches(0.4),  Inches(1.3),  Inches(3.8), Inches(2.2), "📱 Client WhatsApp",
     "Utilisateur final\nvia WhatsApp Business", BLUE),
    (Inches(4.7),  Inches(1.3),  Inches(3.8), Inches(2.2), "🤖 Twilio Studio",
     "Chatbot automatisé\nFlow configurable", RGBColor(0x8B,0x5C,0xF6)),
    (Inches(9.0),  Inches(1.3),  Inches(3.9), Inches(2.2), "💬 Chatwoot",
     "Inbox agent\nSync contacts & messages", RGBColor(0x14,0x7E,0xFF)),
    (Inches(0.4),  Inches(4.0),  Inches(3.8), Inches(2.8), "⚙️ Laravel App",
     "Backend principal\nAPI, Jobs, Webhooks\nScheduler, Queue worker", RED),
    (Inches(4.7),  Inches(4.0),  Inches(3.8), Inches(2.8), "🗄️ MySQL",
     "Base de données\n20+ tables\nConversations, campagnes,\ngames, stats", MID_GRAY),
    (Inches(9.0),  Inches(4.0),  Inches(3.9), Inches(2.8), "🐳 Docker",
     "Déploiement conteneurisé\nPHP 8.3-cli + Node 20\nMigrations auto au boot", GREEN),
]

for lft, top, w, h, title, desc, color in arch_items:
    add_rect(slide, lft, top, w, h, fill_color=WHITE,
             line_color=RGBColor(0xDD,0xDD,0xEE), line_width=Pt(1))
    add_rect(slide, lft, top, w, Inches(0.45), fill_color=color)
    add_text(slide, title, lft + Inches(0.1), top + Inches(0.05), w - Inches(0.2), Inches(0.4),
             font_size=Pt(13), bold=True, color=WHITE)
    add_text(slide, desc, lft + Inches(0.15), top + Inches(0.6), w - Inches(0.3), h - Inches(0.7),
             font_size=Pt(12), color=TEXT_DARK)

# Flèches symboliques
for x in [Inches(4.2), Inches(8.5)]:
    add_text(slide, "⇄", x, Inches(2.0), Inches(0.5), Inches(0.5),
             font_size=Pt(20), bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(slide, "⇅", Inches(1.9), Inches(3.6), Inches(0.5), Inches(0.4),
         font_size=Pt(18), bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(slide, "⇅", Inches(6.2), Inches(3.6), Inches(0.5), Inches(0.4),
         font_size=Pt(18), bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(slide, "⇅", Inches(10.5), Inches(3.6), Inches(0.5), Inches(0.4),
         font_size=Pt(18), bold=True, color=RED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CONCLUSION / MERCI
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_color=DARK_GRAY)
add_rect(slide, 0, 0, Inches(0.35), H, fill_color=RED)
add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), fill_color=RED)

add_text(slide, "WhatsCRM", Inches(0.6), Inches(1.5), Inches(10), Inches(1.2),
         font_size=Pt(58), bold=True, color=WHITE)
add_text(slide, "La solution WhatsApp complète pour TotalEnergies",
         Inches(0.6), Inches(2.8), Inches(10), Inches(0.6),
         font_size=Pt(20), color=ACCENT)

add_rect(slide, Inches(0.6), Inches(3.5), Inches(4.5), Inches(0.04), fill_color=RED)

add_text(slide, "✔  Support client en temps réel\n"
                "✔  Campagnes marketing ciblées\n"
                "✔  Analytics chatbot avancés\n"
                "✔  Gamification & engagement client\n"
                "✔  Déploiement Docker clé en main",
         Inches(0.6), Inches(3.7), Inches(8), Inches(2.8),
         font_size=Pt(16), color=RGBColor(0xCC,0xCC,0xDD))

add_text(slide, "Développé par  YesWeCange  ×  TotalEnergies",
         Inches(0.6), Inches(6.6), Inches(10), Inches(0.5),
         font_size=Pt(13), color=RGBColor(0x77,0x77,0x99), italic=True)


# ─── Sauvegarde ──────────────────────────────────────────────────────────────
output_path = r"C:\wamp64\www\YESWECANGE\totalEnergie-chatbotapp\WhatsCRM_Presentation.pptx"
prs.save(output_path)
print("Fichier genere : " + output_path)
print("Slides : " + str(len(prs.slides)))
